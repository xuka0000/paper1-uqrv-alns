from __future__ import annotations

import math
import shutil
import subprocess
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "_work" / "fig1_ppt_icon_assets_20260613"
OUT_DIR = ROOT / "_work" / "fig1_ppt_icon_outputs_20260613"

SOURCE_DIRS = [
    ROOT
    / "08_rebuild_from_20260606"
    / "paper1_latest_full_package_20260606_highlight_repair_work"
    / "02_latest_manuscript_latex"
    / "published_style_current"
    / "tre_published_style",
    ROOT
    / "08_rebuild_from_20260606"
    / "paper1_latest_full_package_20260606_highlight_repair_work"
    / "02_latest_manuscript_latex"
    / "submission_elsarticle_current"
    / "tre_submission_current",
]


W, H = 960, 620
SCALE = 3

BLACK = (20, 20, 20)
GRAY = (122, 130, 136)
LIGHT_GRAY = (205, 211, 216)
BLUE = (0, 83, 214)
RED = (207, 31, 45)
ORANGE = (242, 128, 27)
GREEN = (37, 139, 55)
PINK = (224, 51, 99)
WHITE = (255, 255, 255)

DEPOT = (80, 342)
RETURN = (880, 342)
STOPS = {
    "S1": (220, 292),
    "S2": (420, 145),
    "S3": (470, 360),
    "S4": (330, 455),
    "S5": (650, 390),
    "S6": (790, 292),
}
TOWERS = {
    "T1": (160, 118, "red"),
    "T2": (250, 112, "orange"),
    "T3": (405, 220, "red"),
    "T4": (470, 200, "red"),
    "T5": (535, 225, "orange"),
    "T6": (420, 58, "green"),
    "T7": (600, 496, "orange"),
    "T8": (690, 496, "red"),
    "T9": (735, 115, "red"),
    "T10": (830, 115, "red"),
}


def sp(v: float) -> int:
    return int(round(v * SCALE))


def recolor_svg(src: Path, dst: Path, color: tuple[int, int, int]) -> None:
    text = src.read_text(encoding="utf-8")
    fill = f'fill="#{color[0]:02x}{color[1]:02x}{color[2]:02x}"'
    if "fill=" in text:
        import re

        text = re.sub(r'fill="[^"]*"', fill, text)
    else:
        text = text.replace("<path ", f"<path {fill} ", 1)
    dst.write_text(text, encoding="utf-8")


def prepare_icons() -> dict[str, Path]:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    icons = {
        "tower_red": ("transmission-tower.svg", RED),
        "tower_orange": ("transmission-tower.svg", ORANGE),
        "tower_green": ("transmission-tower.svg", GREEN),
        "drone_blue": ("quadcopter.svg", BLUE),
        "drone_black": ("quadcopter.svg", BLACK),
        "truck_black": ("truck.svg", BLACK),
        "warehouse_pink": ("warehouse.svg", PINK),
    }
    out: dict[str, Path] = {}
    for key, (name, color) in icons.items():
        svg_path = OUT_DIR / f"{key}.svg"
        png_path = OUT_DIR / f"{key}.png"
        recolor_svg(ASSET_DIR / name, svg_path, color)
        subprocess.run(
            [
                "magick",
                str(svg_path),
                "-background",
                "none",
                "-resize",
                "512x512",
                str(png_path),
            ],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        out[key] = png_path
    return out


def font(size: float, bold: bool = False) -> ImageFont.FreeTypeFont:
    fonts = Path("C:/Windows/Fonts")
    candidates = ["arialbd.ttf" if bold else "arial.ttf", "calibrib.ttf" if bold else "calibri.ttf"]
    for name in candidates:
        path = fonts / name
        if path.exists():
            return ImageFont.truetype(str(path), sp(size))
    return ImageFont.load_default()


def draw_centered(draw: ImageDraw.ImageDraw, text: str, box, size=13, bold=False, fill=BLACK) -> None:
    x, y, w, h = [sp(v) for v in box]
    f = font(size, bold)
    bbox = draw.textbbox((0, 0), text, font=f)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text((x + (w - tw) / 2, y + (h - th) / 2 - sp(0.5)), text, font=f, fill=fill)


def draw_left(draw: ImageDraw.ImageDraw, text: str, box, size=11, bold=False, fill=BLACK) -> None:
    x, y, _w, h = [sp(v) for v in box]
    f = font(size, bold)
    bbox = draw.textbbox((0, 0), text, font=f)
    th = bbox[3] - bbox[1]
    draw.text((x, y + (h - th) / 2), text, font=f, fill=fill)


def paste_icon(canvas: Image.Image, path: Path, x: float, y: float, w: float, h: float) -> None:
    icon = Image.open(path).convert("RGBA").resize((sp(w), sp(h)), Image.LANCZOS)
    canvas.alpha_composite(icon, (sp(x), sp(y)))


def arrow_head(draw: ImageDraw.ImageDraw, p1, p2, color, width=2.0, size=9.0) -> None:
    x1, y1 = p1
    x2, y2 = p2
    angle = math.atan2(y2 - y1, x2 - x1)
    tip = (sp(x2), sp(y2))
    left = (
        sp(x2 - size * math.cos(angle - math.pi / 7)),
        sp(y2 - size * math.sin(angle - math.pi / 7)),
    )
    right = (
        sp(x2 - size * math.cos(angle + math.pi / 7)),
        sp(y2 - size * math.sin(angle + math.pi / 7)),
    )
    draw.polygon([tip, left, right], fill=color)


def draw_line_arrow(draw: ImageDraw.ImageDraw, p1, p2, color=BLACK, width=2.0, dashed=False, arrow=True) -> None:
    x1, y1 = p1
    x2, y2 = p2
    if dashed:
        dash, gap = 10.0, 7.0
        length = math.hypot(x2 - x1, y2 - y1)
        if length == 0:
            return
        ux, uy = (x2 - x1) / length, (y2 - y1) / length
        pos = 0.0
        while pos < length:
            end = min(pos + dash, length)
            a = (sp(x1 + ux * pos), sp(y1 + uy * pos))
            b = (sp(x1 + ux * end), sp(y1 + uy * end))
            draw.line([a, b], fill=color, width=sp(width))
            pos += dash + gap
    else:
        draw.line([(sp(x1), sp(y1)), (sp(x2), sp(y2))], fill=color, width=sp(width))
    if arrow:
        arrow_head(draw, p1, p2, color, width=width)


def draw_stop(draw: ImageDraw.ImageDraw, label: str, x: float, y: float) -> None:
    draw.ellipse([sp(x - 13), sp(y - 13), sp(x + 13), sp(y + 13)], fill=BLACK, outline=BLACK, width=sp(1))
    draw.ellipse([sp(x - 5), sp(y - 5), sp(x + 5), sp(y + 5)], fill=WHITE)
    draw_centered(draw, label, (x - 22, y + 14, 44, 20), size=14, bold=True)


def draw_tower(canvas: Image.Image, draw: ImageDraw.ImageDraw, icons, label: str, x: float, y: float, level: str) -> None:
    paste_icon(canvas, icons[f"tower_{level}"], x - 13, y - 16, 26, 50)
    draw_centered(draw, label, (x - 18, y - 34, 36, 18), size=13, bold=True)


def draw_green_label(draw: ImageDraw.ImageDraw, text: str, x: float, y: float, w: float = 104, h: float = 22) -> None:
    draw.rounded_rectangle(
        [sp(x), sp(y), sp(x + w), sp(y + h)],
        radius=sp(5),
        fill=WHITE,
        outline=GREEN,
        width=sp(1.5),
    )
    draw_centered(draw, text, (x + 2, y + 1, w - 4, h - 2), size=11.5, bold=True, fill=GREEN)


def draw_red_tag(draw: ImageDraw.ImageDraw, text: str, x: float, y: float, w: float = 90, h: float = 22) -> None:
    pts = [
        (sp(x), sp(y)),
        (sp(x + w - 14), sp(y)),
        (sp(x + w), sp(y + h / 2)),
        (sp(x + w - 14), sp(y + h)),
        (sp(x), sp(y + h)),
        (sp(x + 8), sp(y + h / 2)),
    ]
    draw.polygon(pts, fill=RED)
    draw_centered(draw, text, (x + 4, y + 1, w - 17, h - 2), size=11.2, bold=True, fill=WHITE)


def draw_vehicle_link(canvas, draw, icons, a, b, selected=True, truck=True) -> None:
    color = BLACK if selected else LIGHT_GRAY
    width = 2.1 if selected else 1.2
    draw_line_arrow(draw, a, b, color=color, width=width, dashed=not selected)
    if selected and truck:
        mx = (a[0] + b[0]) / 2
        my = (a[1] + b[1]) / 2
        paste_icon(canvas, icons["truck_black"], mx - 10, my - 18, 20, 20)


def draw_uav_sortie(canvas, draw, icons, points, drone_at=None) -> None:
    for a, b in zip(points, points[1:]):
        draw_line_arrow(draw, a, b, color=BLUE, width=1.65, dashed=True)
    if drone_at:
        paste_icon(canvas, icons["drone_black"], drone_at[0] - 9, drone_at[1] - 8, 20, 18)


def draw_legend(canvas, draw, icons) -> None:
    y = 542
    draw.rounded_rectangle([sp(45), sp(y), sp(915), sp(y + 55)], radius=sp(6), fill=WHITE, outline=(80, 80, 80), width=sp(1.1))
    x = 65
    draw.ellipse([sp(x), sp(y + 18), sp(x + 18), sp(y + 36)], fill=BLACK)
    draw.ellipse([sp(x + 6), sp(y + 24), sp(x + 12), sp(y + 30)], fill=WHITE)
    draw_left(draw, "candidate stop", (x + 24, y + 14, 92, 26), size=10.8)

    x = 205
    paste_icon(canvas, icons["tower_red"], x, y + 9, 17, 35)
    paste_icon(canvas, icons["tower_orange"], x + 25, y + 9, 17, 35)
    paste_icon(canvas, icons["tower_green"], x + 50, y + 9, 17, 35)
    draw_left(draw, "tower priority", (x + 75, y + 14, 83, 26), size=10.8)

    x = 370
    draw_line_arrow(draw, (x, y + 28), (x + 55, y + 28), color=BLACK, width=1.8)
    draw_left(draw, "vehicle route", (x + 64, y + 14, 80, 26), size=10.8)

    x = 535
    draw_line_arrow(draw, (x, y + 28), (x + 55, y + 28), color=BLUE, width=1.5, dashed=True)
    draw_left(draw, "same-stop UAV sortie", (x + 64, y + 14, 128, 26), size=10.8)

    x = 730
    draw.rounded_rectangle([sp(x), sp(y + 18), sp(x + 42), sp(y + 38)], radius=sp(4), fill=WHITE, outline=GREEN, width=sp(1.2))
    draw_left(draw, "q95 feasible sortie", (x + 50, y + 14, 116, 26), size=10.8)


def render_png(icons: dict[str, Path]) -> Path:
    canvas = Image.new("RGBA", (sp(W), sp(H)), WHITE + (255,))
    draw = ImageDraw.Draw(canvas)

    for a, b in [
        (STOPS["S1"], STOPS["S2"]),
        (STOPS["S2"], STOPS["S5"]),
        (STOPS["S1"], STOPS["S4"]),
        (STOPS["S4"], STOPS["S3"]),
    ]:
        draw_vehicle_link(canvas, draw, icons, a, b, selected=False, truck=False)
    for a, b in [
        (DEPOT, STOPS["S1"]),
        (STOPS["S1"], STOPS["S3"]),
        (STOPS["S3"], STOPS["S5"]),
        (STOPS["S5"], STOPS["S6"]),
        (STOPS["S6"], RETURN),
    ]:
        draw_vehicle_link(canvas, draw, icons, a, b, selected=True, truck=True)

    paste_icon(canvas, icons["warehouse_pink"], DEPOT[0] - 34, DEPOT[1] - 36, 54, 54)
    draw_centered(draw, "Depot", (DEPOT[0] - 40, DEPOT[1] + 20, 80, 22), size=14, bold=True)
    paste_icon(canvas, icons["warehouse_pink"], RETURN[0] - 16, RETURN[1] - 36, 54, 54)
    draw_centered(draw, "Return depot", (RETURN[0] - 43, RETURN[1] + 20, 96, 22), size=14, bold=True)

    for label, (x, y) in STOPS.items():
        draw_stop(draw, label, x, y)
    for label, (x, y, level) in TOWERS.items():
        draw_tower(canvas, draw, icons, label, x, y, level)

    draw_uav_sortie(canvas, draw, icons, [STOPS["S1"], (160, 152), (250, 146), STOPS["S1"]], drone_at=(207, 160))
    draw_uav_sortie(canvas, draw, icons, [STOPS["S3"], (405, 253), (470, 235), (535, 258), STOPS["S3"]], drone_at=(545, 278))
    draw_uav_sortie(canvas, draw, icons, [STOPS["S5"], (600, 526), STOPS["S5"]], drone_at=(623, 520))
    draw_uav_sortie(canvas, draw, icons, [STOPS["S5"], (690, 526), STOPS["S5"]], drone_at=(675, 520))
    draw_uav_sortie(canvas, draw, icons, [STOPS["S6"], (735, 150), (830, 150), STOPS["S6"]], drone_at=(782, 162))
    draw_uav_sortie(canvas, draw, icons, [STOPS["S2"], (420, 92), STOPS["S2"]], drone_at=(446, 101))

    draw_green_label(draw, "q95 feasible", 455, 52, 82, 22)

    draw_red_tag(draw, "T1 T2 high", 112, 182, 94, 22)
    draw_green_label(draw, "S1: T1 T2", 252, 198, 86, 22)
    draw_red_tag(draw, "T3 T4 T5", 395, 285, 104, 22)
    draw_green_label(draw, "S3: T3-T5", 510, 315, 96, 22)
    draw_green_label(draw, "parallel UAVs", 675, 442, 102, 22)
    draw_red_tag(draw, "T8 high", 712, 500, 76, 22)
    draw_red_tag(draw, "T9 T10 high", 775, 188, 106, 22)
    draw_green_label(draw, "S6: T9-T10", 812, 222, 96, 22)

    draw_legend(canvas, draw, icons)

    png = OUT_DIR / "Fig0_network_construction_ppt_icon.png"
    pdf = OUT_DIR / "Fig0_network_construction_ppt_icon.pdf"
    canvas.convert("RGB").save(png, quality=95)
    canvas.convert("RGB").save(pdf, "PDF", resolution=300.0)
    return png


def ppt_rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(color[0], color[1], color[2])


def pt(v: float):
    return Pt(v)


def add_ppt_text(slide, text, x, y, w, h, size=13, bold=False, color=BLACK, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(pt(x), pt(y), pt(w), pt(h))
    tf = box.text_frame
    tf.margin_left = pt(1)
    tf.margin_right = pt(1)
    tf.margin_top = pt(1)
    tf.margin_bottom = pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = pt(size)
    run.font.bold = bold
    run.font.color.rgb = ppt_rgb(color)
    return box


def add_ppt_line(slide, p1, p2, color=BLACK, width=2.0, dashed=False, arrow=True):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, pt(p1[0]), pt(p1[1]), pt(p2[0]), pt(p2[1]))
    line.line.color.rgb = ppt_rgb(color)
    line.line.width = pt(width)
    if dashed:
        line.line.dash_style = 4
    if arrow:
        line._element.spPr.ln.append(
            parse_xml(f'<a:tailEnd {nsdecls("a")} type="triangle" w="med" len="med"/>')
        )
    return line


def add_ppt_stop(slide, label, x, y):
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, pt(x - 13), pt(y - 13), pt(26), pt(26))
    outer.fill.solid()
    outer.fill.fore_color.rgb = ppt_rgb(BLACK)
    outer.line.color.rgb = ppt_rgb(BLACK)
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, pt(x - 5), pt(y - 5), pt(10), pt(10))
    inner.fill.solid()
    inner.fill.fore_color.rgb = ppt_rgb(WHITE)
    inner.line.fill.background()
    add_ppt_text(slide, label, x - 22, y + 14, 44, 20, size=14, bold=True)


def add_ppt_icon(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), pt(x), pt(y), width=pt(w), height=pt(h))


def add_ppt_green_label(slide, text, x, y, w=104, h=22):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pt(x), pt(y), pt(w), pt(h))
    box.fill.solid()
    box.fill.fore_color.rgb = ppt_rgb(WHITE)
    box.line.color.rgb = ppt_rgb(GREEN)
    box.line.width = pt(1.2)
    add_ppt_text(slide, text, x + 2, y + 1, w - 4, h - 2, size=11.5, bold=True, color=GREEN)


def add_ppt_red_tag(slide, text, x, y, w=90, h=22):
    tag = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, pt(x), pt(y), pt(w), pt(h))
    tag.fill.solid()
    tag.fill.fore_color.rgb = ppt_rgb(RED)
    tag.line.color.rgb = ppt_rgb(RED)
    add_ppt_text(slide, text, x + 4, y + 1, w - 17, h - 2, size=11.2, bold=True, color=WHITE)


def build_pptx(icons: dict[str, Path]) -> Path:
    prs = Presentation()
    prs.slide_width = pt(W)
    prs.slide_height = pt(H)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    for a, b in [
        (STOPS["S1"], STOPS["S2"]),
        (STOPS["S2"], STOPS["S5"]),
        (STOPS["S1"], STOPS["S4"]),
        (STOPS["S4"], STOPS["S3"]),
    ]:
        add_ppt_line(slide, a, b, color=LIGHT_GRAY, width=1.2, dashed=True)
    for a, b in [
        (DEPOT, STOPS["S1"]),
        (STOPS["S1"], STOPS["S3"]),
        (STOPS["S3"], STOPS["S5"]),
        (STOPS["S5"], STOPS["S6"]),
        (STOPS["S6"], RETURN),
    ]:
        add_ppt_line(slide, a, b, color=BLACK, width=2.1)
        mx, my = (a[0] + b[0]) / 2, (a[1] + b[1]) / 2
        add_ppt_icon(slide, icons["truck_black"], mx - 10, my - 18, 20, 20)

    add_ppt_icon(slide, icons["warehouse_pink"], DEPOT[0] - 34, DEPOT[1] - 36, 54, 54)
    add_ppt_text(slide, "Depot", DEPOT[0] - 40, DEPOT[1] + 20, 80, 22, size=14, bold=True)
    add_ppt_icon(slide, icons["warehouse_pink"], RETURN[0] - 16, RETURN[1] - 36, 54, 54)
    add_ppt_text(slide, "Return depot", RETURN[0] - 43, RETURN[1] + 20, 96, 22, size=14, bold=True)

    for label, (x, y) in STOPS.items():
        add_ppt_stop(slide, label, x, y)
    for label, (x, y, level) in TOWERS.items():
        add_ppt_icon(slide, icons[f"tower_{level}"], x - 13, y - 16, 26, 50)
        add_ppt_text(slide, label, x - 18, y - 34, 36, 18, size=13, bold=True)

    sorties = [
        ([STOPS["S1"], (160, 152), (250, 146), STOPS["S1"]], (207, 160)),
        ([STOPS["S3"], (405, 253), (470, 235), (535, 258), STOPS["S3"]], (545, 278)),
        ([STOPS["S5"], (600, 526), STOPS["S5"]], (623, 520)),
        ([STOPS["S5"], (690, 526), STOPS["S5"]], (675, 520)),
        ([STOPS["S6"], (735, 150), (830, 150), STOPS["S6"]], (782, 162)),
        ([STOPS["S2"], (420, 92), STOPS["S2"]], (446, 101)),
    ]
    for pts, drone_at in sorties:
        for a, b in zip(pts, pts[1:]):
            add_ppt_line(slide, a, b, color=BLUE, width=1.65, dashed=True)
        add_ppt_icon(slide, icons["drone_black"], drone_at[0] - 9, drone_at[1] - 8, 20, 18)

    add_ppt_green_label(slide, "q95 feasible", 455, 52, 82, 22)

    add_ppt_red_tag(slide, "T1 T2 high", 112, 182, 94, 22)
    add_ppt_green_label(slide, "S1: T1 T2", 252, 198, 86, 22)
    add_ppt_red_tag(slide, "T3 T4 T5", 395, 285, 104, 22)
    add_ppt_green_label(slide, "S3: T3-T5", 510, 315, 96, 22)
    add_ppt_green_label(slide, "parallel UAVs", 675, 442, 102, 22)
    add_ppt_red_tag(slide, "T8 high", 712, 500, 76, 22)
    add_ppt_red_tag(slide, "T9 T10 high", 775, 188, 106, 22)
    add_ppt_green_label(slide, "S6: T9-T10", 812, 222, 96, 22)

    # Legend source elements remain editable in PowerPoint.
    leg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pt(45), pt(542), pt(870), pt(55))
    leg.fill.solid()
    leg.fill.fore_color.rgb = ppt_rgb(WHITE)
    leg.line.color.rgb = ppt_rgb((80, 80, 80))
    add_ppt_stop(slide, "", 74, 569)
    add_ppt_text(slide, "candidate stop", 89, 556, 92, 26, size=10.8, align=PP_ALIGN.LEFT)
    add_ppt_icon(slide, icons["tower_red"], 205, 551, 17, 35)
    add_ppt_icon(slide, icons["tower_orange"], 230, 551, 17, 35)
    add_ppt_icon(slide, icons["tower_green"], 255, 551, 17, 35)
    add_ppt_text(slide, "tower priority", 280, 556, 83, 26, size=10.8, align=PP_ALIGN.LEFT)
    add_ppt_line(slide, (370, 570), (425, 570), color=BLACK, width=1.8)
    add_ppt_text(slide, "vehicle route", 434, 556, 80, 26, size=10.8, align=PP_ALIGN.LEFT)
    add_ppt_line(slide, (535, 570), (590, 570), color=BLUE, width=1.5, dashed=True)
    add_ppt_text(slide, "same-stop UAV sortie", 599, 556, 128, 26, size=10.8, align=PP_ALIGN.LEFT)
    add_ppt_green_label(slide, "", 730, 560, 42, 20)
    add_ppt_text(slide, "q95 feasible sortie", 780, 556, 116, 26, size=10.8, align=PP_ALIGN.LEFT)

    pptx = OUT_DIR / "Fig0_network_construction_ppt_icon_source.pptx"
    prs.save(pptx)
    return pptx


def install_outputs(pptx: Path, png: Path, pdf: Path) -> None:
    for src_dir in SOURCE_DIRS:
        fig_dir = src_dir / "figures"
        fig_dir.mkdir(parents=True, exist_ok=True)
        shutil.copy2(png, fig_dir / "Fig0_network_construction_ai.png")
        shutil.copy2(png, fig_dir / "Fig0_network_construction_ppt_icon.png")
        shutil.copy2(pdf, fig_dir / "Fig0_network_construction_ppt_icon.pdf")
        shutil.copy2(pptx, fig_dir / "Fig0_network_construction_ppt_icon_source.pptx")


if __name__ == "__main__":
    icons = prepare_icons()
    pptx_path = build_pptx(icons)
    png_path = render_png(icons)
    pdf_path = OUT_DIR / "Fig0_network_construction_ppt_icon.pdf"
    install_outputs(pptx_path, png_path, pdf_path)
    print(f"PPTX={pptx_path}")
    print(f"PNG={png_path}")
    print(f"PDF={pdf_path}")
    for src in SOURCE_DIRS:
        print(f"INSTALLED={src / 'figures' / 'Fig0_network_construction_ai.png'}")
