from __future__ import annotations

import importlib.util
import shutil
import subprocess
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


ROOT = Path(__file__).resolve().parents[1]
BASE_SCRIPT = ROOT / "scripts" / "build_fig1_ppt_icon_network.py"
SPEC = importlib.util.spec_from_file_location("fig_base", BASE_SCRIPT)
base = importlib.util.module_from_spec(SPEC)
assert SPEC.loader is not None
SPEC.loader.exec_module(base)

OUT_DIR = ROOT / "_work" / "reference_style_icon_figures_20260613"
base.OUT_DIR = OUT_DIR

SOURCE_DIRS = [
    ROOT
    / "08_rebuild_from_20260606"
    / "paper1_latest_full_package_20260606_highlight_repair_work"
    / "02_latest_manuscript_latex"
    / "published_style_current"
    / "tre_published_style",
    ROOT
    / "08_rebuild_from_20260606"
    / "paper1_latest_full_package_20260606_highlight_repair_work"
    / "02_latest_manuscript_latex"
    / "submission_elsarticle_current"
    / "tre_submission_current",
]


W = 960
H = 650


def new_canvas(h: int = H) -> tuple[Image.Image, ImageDraw.ImageDraw]:
    canvas = Image.new("RGBA", (base.sp(W), base.sp(h)), base.WHITE + (255,))
    return canvas, ImageDraw.Draw(canvas)


def save_canvas(canvas: Image.Image, stem: str) -> tuple[Path, Path]:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    png = OUT_DIR / f"{stem}.png"
    pdf = OUT_DIR / f"{stem}.pdf"
    canvas.convert("RGB").save(png, quality=95)
    canvas.convert("RGB").save(pdf, "PDF", resolution=300.0)
    return png, pdf


def new_ppt(h: int = H):
    prs = Presentation()
    prs.slide_width = base.pt(W)
    prs.slide_height = base.pt(h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def extend_icons(icons: dict[str, Path]) -> dict[str, Path]:
    extra = {
        "tower_blue": ((0, 102, 204), "transmission-tower.svg"),
        "tower_purple": ((100, 45, 160), "transmission-tower.svg"),
        "tower_gold": ((213, 150, 0), "transmission-tower.svg"),
        "tower_gray": ((160, 160, 160), "transmission-tower.svg"),
    }
    for key, (color, src_name) in extra.items():
        svg_path = OUT_DIR / f"{key}.svg"
        png_path = OUT_DIR / f"{key}.png"
        base.recolor_svg(base.ASSET_DIR / src_name, svg_path, color)
        subprocess.run(
            ["magick", str(svg_path), "-background", "none", "-resize", "512x512", str(png_path)],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        icons[key] = png_path
    return icons


def ppt_save(prs, stem: str) -> Path:
    pptx = OUT_DIR / f"{stem}_source.pptx"
    prs.save(pptx)
    return pptx


def dashed_box(draw, x, y, w, h, color=base.BLACK, width=1.5):
    dash = 10
    gap = 7
    pts = [(x, y, x + w, y), (x + w, y, x + w, y + h), (x + w, y + h, x, y + h), (x, y + h, x, y)]
    for x1, y1, x2, y2 in pts:
        length = ((x2 - x1) ** 2 + (y2 - y1) ** 2) ** 0.5
        if length == 0:
            continue
        ux = (x2 - x1) / length
        uy = (y2 - y1) / length
        pos = 0
        while pos < length:
            end = min(pos + dash, length)
            draw.line(
                [
                    (base.sp(x1 + ux * pos), base.sp(y1 + uy * pos)),
                    (base.sp(x1 + ux * end), base.sp(y1 + uy * end)),
                ],
                fill=color,
                width=base.sp(width),
            )
            pos += dash + gap


def ppt_dashed_box(slide, x, y, w, h, color=base.BLACK):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, base.pt(x), base.pt(y), base.pt(w), base.pt(h))
    rect.fill.background()
    rect.line.color.rgb = base.ppt_rgb(color)
    rect.line.width = base.pt(1.4)
    rect.line.dash_style = 4
    return rect


def green_tag(draw, text, x, y, w=92, h=22):
    base.draw_green_label(draw, text, x, y, w, h)


def red_tag(draw, text, x, y, w=92, h=22):
    base.draw_red_tag(draw, text, x, y, w, h)


def blue_tag(draw, text, x, y, w=88, h=22):
    draw = draw
    draw.rounded_rectangle(
        [base.sp(x), base.sp(y), base.sp(x + w), base.sp(y + h)],
        radius=base.sp(5),
        fill=base.WHITE,
        outline=base.BLUE,
        width=base.sp(1.5),
    )
    base.draw_centered(draw, text, (x + 2, y + 1, w - 4, h - 2), size=11.2, bold=True, fill=base.BLUE)


def mini_tower(canvas, draw, icons, label, x, y, level, icon_w=18, icon_h=34):
    base.paste_icon(canvas, icons[f"tower_{level}"], x - icon_w / 2, y, icon_w, icon_h)
    base.draw_centered(draw, label, (x - 18, y - 16, 36, 14), size=9.8, bold=True)


def ppt_blue_tag(slide, text, x, y, w=88, h=22):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, base.pt(x), base.pt(y), base.pt(w), base.pt(h))
    box.fill.solid()
    box.fill.fore_color.rgb = base.ppt_rgb(base.WHITE)
    box.line.color.rgb = base.ppt_rgb(base.BLUE)
    box.line.width = base.pt(1.2)
    base.add_ppt_text(slide, text, x + 2, y + 1, w - 4, h - 2, size=11.2, bold=True, color=base.BLUE)


def ppt_mini_tower(slide, icons, label, x, y, level, icon_w=18, icon_h=34):
    base.add_ppt_icon(slide, icons[f"tower_{level}"], x - icon_w / 2, y, icon_w, icon_h)
    base.add_ppt_text(slide, label, x - 18, y - 16, 36, 14, size=9.8, bold=True)


def draw_tower_row(canvas, draw, icons, towers, y_label=True):
    for label, x, y, level in towers:
        base.draw_tower(canvas, draw, icons, label, x, y, level)


def draw_fig1(icons):
    h = 670
    canvas, draw = new_canvas(h)
    depot = (70, 330)
    ret = (890, 330)
    stops = {
        "S1": (210, 275),
        "S2": (410, 145),
        "S3": (470, 355),
        "S4": (310, 452),
        "S5": (665, 382),
        "S6": (795, 275),
    }
    towers = {
        "T1": (145, 105, "red"),
        "T2": (240, 98, "orange"),
        "T3": (395, 220, "red"),
        "T4": (470, 195, "red"),
        "T5": (545, 222, "orange"),
        "T6": (410, 55, "green"),
        "T7": (602, 492, "orange"),
        "T8": (700, 492, "red"),
        "T9": (735, 100, "orange"),
        "T10": (835, 100, "red"),
    }
    for a, b in [
        (stops["S1"], stops["S2"]),
        (stops["S2"], stops["S5"]),
        (stops["S1"], stops["S4"]),
        (stops["S4"], stops["S3"]),
    ]:
        base.draw_vehicle_link(canvas, draw, icons, a, b, selected=False, truck=False)
    for a, b in [
        (depot, stops["S1"]),
        (stops["S1"], stops["S3"]),
        (stops["S3"], stops["S5"]),
        (stops["S5"], stops["S6"]),
        (stops["S6"], ret),
    ]:
        base.draw_vehicle_link(canvas, draw, icons, a, b, selected=True, truck=True)
    base.paste_icon(canvas, icons["warehouse_pink"], depot[0] - 35, depot[1] - 35, 54, 54)
    base.draw_centered(draw, "Depot", (depot[0] - 45, depot[1] + 20, 80, 22), size=14, bold=True)
    base.paste_icon(canvas, icons["warehouse_pink"], ret[0] - 18, ret[1] - 35, 54, 54)
    base.draw_centered(draw, "Return depot", (ret[0] - 48, ret[1] + 20, 104, 22), size=14, bold=True)
    for label, (x, y) in stops.items():
        base.draw_stop(draw, label, x, y)
    for label, (x, y, level) in towers.items():
        base.draw_tower(canvas, draw, icons, label, x, y, level)
    sorties = [
        ([stops["S1"], (145, 140), (240, 136), stops["S1"]], (198, 150)),
        ([stops["S3"], (395, 255), (470, 230), (545, 258), stops["S3"]], (548, 278)),
        ([stops["S5"], (602, 522), stops["S5"]], (625, 520)),
        ([stops["S5"], (700, 522), stops["S5"]], (675, 520)),
        ([stops["S6"], (735, 134), (835, 134), stops["S6"]], (788, 150)),
        ([stops["S2"], (410, 90), stops["S2"]], (438, 102)),
    ]
    for pts, drone_at in sorties:
        base.draw_uav_sortie(canvas, draw, icons, pts, drone_at=drone_at)
    red_tag(draw, "P1: T1 T2", 105, 170, 90)
    green_tag(draw, "S1: P1", 255, 185, 76)
    red_tag(draw, "P2: T3 T4 T5", 395, 290, 112)
    green_tag(draw, "S3: P2", 515, 318, 80)
    green_tag(draw, "q95 feasible", 445, 52, 88)
    red_tag(draw, "P3: T7 T8", 715, 505, 92)
    green_tag(draw, "parallel UAVs", 675, 435, 102)
    red_tag(draw, "P4: T9 T10", 772, 175, 102)
    green_tag(draw, "S6: P4", 815, 210, 78)
    green_tag(draw, "reserve ok", 278, 215, 82)
    green_tag(draw, "q95 service", 570, 340, 88)
    base.draw_legend(canvas, draw, icons)
    png, pdf = save_canvas(canvas, "Fig0_network_construction_ppt_icon")

    prs, slide = new_ppt(h)
    # Editable source mirrors the same element grammar, with the rendered PNG as a locked background reference removed from final rendering.
    # The source remains editable because all meaningful labels, icons and connectors are recreated as native objects.
    for a, b in [
        (stops["S1"], stops["S2"]),
        (stops["S2"], stops["S5"]),
        (stops["S1"], stops["S4"]),
        (stops["S4"], stops["S3"]),
    ]:
        base.add_ppt_line(slide, a, b, color=base.LIGHT_GRAY, width=1.2, dashed=True)
    for a, b in [(depot, stops["S1"]), (stops["S1"], stops["S3"]), (stops["S3"], stops["S5"]), (stops["S5"], stops["S6"]), (stops["S6"], ret)]:
        base.add_ppt_line(slide, a, b, color=base.BLACK, width=2.1)
        mx, my = (a[0] + b[0]) / 2, (a[1] + b[1]) / 2
        base.add_ppt_icon(slide, icons["truck_black"], mx - 10, my - 18, 20, 20)
    base.add_ppt_icon(slide, icons["warehouse_pink"], depot[0] - 35, depot[1] - 35, 54, 54)
    base.add_ppt_text(slide, "Depot", depot[0] - 45, depot[1] + 20, 80, 22, size=14, bold=True)
    base.add_ppt_icon(slide, icons["warehouse_pink"], ret[0] - 18, ret[1] - 35, 54, 54)
    base.add_ppt_text(slide, "Return depot", ret[0] - 48, ret[1] + 20, 104, 22, size=14, bold=True)
    for label, (x, y) in stops.items():
        base.add_ppt_stop(slide, label, x, y)
    for label, (x, y, level) in towers.items():
        base.add_ppt_icon(slide, icons[f"tower_{level}"], x - 13, y - 16, 26, 50)
        base.add_ppt_text(slide, label, x - 18, y - 34, 36, 18, size=13, bold=True)
    for pts, drone_at in sorties:
        for a, b in zip(pts, pts[1:]):
            base.add_ppt_line(slide, a, b, color=base.BLUE, width=1.65, dashed=True)
        base.add_ppt_icon(slide, icons["drone_black"], drone_at[0] - 9, drone_at[1] - 8, 20, 18)
    for text, x, y, w in [
        ("P1: T1 T2", 105, 170, 90),
        ("P2: T3 T4 T5", 395, 290, 112),
        ("P3: T7 T8", 715, 505, 92),
        ("P4: T9 T10", 772, 175, 102),
    ]:
        base.add_ppt_red_tag(slide, text, x, y, w)
    for text, x, y, w in [
        ("S1: P1", 255, 185, 76),
        ("S3: P2", 515, 318, 80),
        ("q95 feasible", 445, 52, 88),
        ("parallel UAVs", 675, 435, 102),
        ("S6: P4", 815, 210, 78),
        ("reserve ok", 278, 215, 82),
        ("q95 service", 570, 340, 88),
    ]:
        base.add_ppt_green_label(slide, text, x, y, w, 22)
    pptx = ppt_save(prs, "Fig0_network_construction_ppt_icon")
    return png, pdf, pptx


def draw_fig2(icons):
    h = 650
    canvas, draw = new_canvas(h)
    base.draw_left(draw, "Current solution:", (28, 22, 160, 24), size=15, bold=True)
    dep = (78, 135)
    ret = (884, 135)
    stops = {"S1": (225, 135), "S3": (435, 135), "S5": (640, 135), "S6": (795, 135)}
    for a, b in [(dep, stops["S1"]), (stops["S1"], stops["S3"]), (stops["S3"], stops["S5"]), (stops["S5"], stops["S6"]), (stops["S6"], ret)]:
        base.draw_vehicle_link(canvas, draw, icons, a, b, selected=True, truck=True)
    base.paste_icon(canvas, icons["warehouse_pink"], dep[0] - 28, dep[1] - 30, 46, 46)
    base.draw_centered(draw, "Depot", (dep[0] - 35, dep[1] + 22, 70, 18), size=12, bold=True)
    base.paste_icon(canvas, icons["warehouse_pink"], ret[0] - 18, ret[1] - 30, 46, 46)
    base.draw_centered(draw, "Return", (ret[0] - 35, ret[1] + 22, 70, 18), size=12, bold=True)
    for s, p in stops.items():
        base.draw_stop(draw, s, *p)
    groups = [
        ("S1", [("T1", 190, 50, "green"), ("T2", 255, 50, "orange")], "P1: T1 T2", 172, 192),
        ("S3", [("T3", 380, 45, "red"), ("T4", 435, 45, "orange"), ("T5", 490, 45, "red")], "P2: T3 T4 T5", 380, 192),
        ("S5", [("T7", 610, 50, "orange"), ("T8", 675, 50, "red")], "P3: T7 T8", 592, 192),
        ("S6", [("T9", 765, 50, "orange"), ("T10", 835, 50, "red")], "P4: T9 T10", 755, 192),
    ]
    for s, tws, tag, tx, ty in groups:
        stop = stops[s]
        pts = [stop] + [(x, y + 32) for _label, x, y, _level in tws] + [stop]
        for a, b in zip(pts, pts[1:]):
            base.draw_line_arrow(draw, a, b, color=base.BLUE, width=1.4, dashed=True)
        base.paste_icon(canvas, icons["drone_black"], stop[0] + 25, stop[1] - 60, 18, 16)
        for label, x, y, level in tws:
            base.draw_tower(canvas, draw, icons, label, x, y, level)
        red_tag(draw, tag, tx, ty, 112 if "T3" in tag else 92)

    base.draw_left(draw, "Candidate service list:", (28, 254, 220, 24), size=15, bold=True)
    dashed_box(draw, 58, 295, 365, 260)
    rows = [
        ("S1", [("T1", "green"), ("T2", "orange")], "P1: T1 T2"),
        ("S3", [("T3", "red"), ("T4", "orange"), ("T5", "red")], "P2: T3 T4 T5"),
        ("S5", [("T7", "orange"), ("T8", "red")], "P3: T7 T8"),
        ("S6", [("T9", "orange"), ("T10", "red")], "P4: T9 T10"),
    ]
    y0 = 318
    for idx, (s, tws, tag) in enumerate(rows):
        y = y0 + idx * 56
        base.draw_stop(draw, "", 92, y + 20)
        base.draw_left(draw, s, (114, y + 8, 38, 24), size=13, bold=True)
        for j, (tw, lev) in enumerate(tws):
            x = 178 + j * 38
            base.draw_tower(canvas, draw, icons, tw, x, y + 2, lev)
        red_tag(draw, tag, 300, y + 10, 108 if "T3" in tag else 88)
    base.draw_centered(draw, "Splitting", (452, 355, 90, 24), size=14, bold=True)
    base.draw_line_arrow(draw, (440, 392), (545, 392), color=base.BLACK, width=5)
    green_tag(draw, "q95 screen", 462, 430, 105, 24)
    base.paste_icon(canvas, icons["drone_black"], 505, 460, 28, 24)
    dashed_box(draw, 585, 295, 335, 260)
    split_rows = [
        ("S1", ["T1", "T2"], "q95 feasible"),
        ("S3", ["T3", "T4"], "q95 feasible"),
        ("S3", ["T5"], "q95 feasible"),
        ("S5", ["T7", "T8"], "parallel ok"),
        ("S6", ["T9", "T10"], "q95 feasible"),
    ]
    for idx, (s, tws, q) in enumerate(split_rows):
        y = 312 + idx * 45
        base.draw_stop(draw, "", 612, y + 20)
        base.draw_left(draw, s, (633, y + 8, 36, 24), size=12.5, bold=True)
        for j, tw in enumerate(tws):
            level = {"T1": "green", "T2": "orange", "T3": "red", "T4": "orange", "T5": "red", "T7": "orange", "T8": "red", "T9": "orange", "T10": "red"}[tw]
            mini_tower(canvas, draw, icons, tw, 705 + j * 42, y + 6, level)
        blue_tag(draw, q, 805, y + 14, 92)

    # Legend
    ly = 590
    draw.rounded_rectangle([base.sp(45), base.sp(ly), base.sp(915), base.sp(ly + 42)], radius=base.sp(6), fill=base.WHITE, outline=(80, 80, 80), width=base.sp(1.1))
    red_tag(draw, "grouped pattern", 70, ly + 10, 122, 20)
    blue_tag(draw, "feasible subpattern", 250, ly + 10, 128, 20)
    base.draw_line_arrow(draw, (430, ly + 22), (490, ly + 22), color=base.BLACK, width=1.7)
    base.draw_left(draw, "vehicle route", (500, ly + 9, 90, 24), size=10.5)
    base.draw_line_arrow(draw, (620, ly + 22), (680, ly + 22), color=base.BLUE, width=1.5, dashed=True)
    base.draw_left(draw, "same-stop UAV sortie", (690, ly + 9, 140, 24), size=10.5)
    green_tag(draw, "q95 screen", 830, ly + 10, 78, 20)
    png, pdf = save_canvas(canvas, "Fig2_solution_framework_ppt_icon")

    prs, slide = new_ppt(h)
    base.add_ppt_text(slide, "Current solution:", 28, 22, 160, 24, size=15, bold=True, align=PP_ALIGN.LEFT)
    for a, b in [(dep, stops["S1"]), (stops["S1"], stops["S3"]), (stops["S3"], stops["S5"]), (stops["S5"], stops["S6"]), (stops["S6"], ret)]:
        base.add_ppt_line(slide, a, b, color=base.BLACK, width=2.0)
        mx, my = (a[0] + b[0]) / 2, (a[1] + b[1]) / 2
        base.add_ppt_icon(slide, icons["truck_black"], mx - 10, my - 18, 20, 20)
    base.add_ppt_icon(slide, icons["warehouse_pink"], dep[0] - 28, dep[1] - 30, 46, 46)
    base.add_ppt_text(slide, "Depot", dep[0] - 35, dep[1] + 22, 70, 18, size=12, bold=True)
    base.add_ppt_icon(slide, icons["warehouse_pink"], ret[0] - 18, ret[1] - 30, 46, 46)
    base.add_ppt_text(slide, "Return", ret[0] - 35, ret[1] + 22, 70, 18, size=12, bold=True)
    for s, p in stops.items():
        base.add_ppt_stop(slide, s, *p)
    for s, tws, tag, tx, ty in groups:
        stop = stops[s]
        pts = [stop] + [(x, y + 32) for _label, x, y, _level in tws] + [stop]
        for a, b in zip(pts, pts[1:]):
            base.add_ppt_line(slide, a, b, color=base.BLUE, width=1.4, dashed=True)
        base.add_ppt_icon(slide, icons["drone_black"], stop[0] + 25, stop[1] - 60, 18, 16)
        for label, x, y, level in tws:
            base.add_ppt_icon(slide, icons[f"tower_{level}"], x - 13, y - 16, 26, 50)
            base.add_ppt_text(slide, label, x - 18, y - 34, 36, 18, size=13, bold=True)
        base.add_ppt_red_tag(slide, tag, tx, ty, 112 if "T3" in tag else 92, 22)

    base.add_ppt_text(slide, "Candidate service list:", 28, 254, 220, 24, size=15, bold=True, align=PP_ALIGN.LEFT)
    ppt_dashed_box(slide, 58, 295, 365, 260)
    for idx, (s, tws, tag) in enumerate(rows):
        y = y0 + idx * 56
        base.add_ppt_stop(slide, "", 92, y + 20)
        base.add_ppt_text(slide, s, 114, y + 8, 38, 24, size=13, bold=True, align=PP_ALIGN.LEFT)
        for j, (tw, lev) in enumerate(tws):
            x = 178 + j * 38
            base.add_ppt_icon(slide, icons[f"tower_{lev}"], x - 13, y - 14, 26, 50)
            base.add_ppt_text(slide, tw, x - 18, y - 32, 36, 18, size=12, bold=True)
        base.add_ppt_red_tag(slide, tag, 300, y + 10, 108 if "T3" in tag else 88, 22)
    base.add_ppt_text(slide, "Splitting", 452, 355, 90, 24, size=14, bold=True)
    base.add_ppt_line(slide, (440, 392), (545, 392), color=base.BLACK, width=5)
    base.add_ppt_green_label(slide, "q95 screen", 462, 430, 105, 24)
    base.add_ppt_icon(slide, icons["drone_black"], 505, 460, 28, 24)
    ppt_dashed_box(slide, 585, 295, 335, 260)
    for idx, (s, tws, q) in enumerate(split_rows):
        y = 312 + idx * 45
        base.add_ppt_stop(slide, "", 612, y + 20)
        base.add_ppt_text(slide, s, 633, y + 8, 36, 24, size=12.5, bold=True, align=PP_ALIGN.LEFT)
        for j, tw in enumerate(tws):
            level = {"T1": "green", "T2": "orange", "T3": "red", "T4": "orange", "T5": "red", "T7": "orange", "T8": "red", "T9": "orange", "T10": "red"}[tw]
            ppt_mini_tower(slide, icons, tw, 705 + j * 42, y + 6, level)
        ppt_blue_tag(slide, q, 805, y + 14, 92)
    leg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, base.pt(45), base.pt(590), base.pt(870), base.pt(42))
    leg.fill.solid()
    leg.fill.fore_color.rgb = base.ppt_rgb(base.WHITE)
    leg.line.color.rgb = base.ppt_rgb((80, 80, 80))
    base.add_ppt_red_tag(slide, "grouped pattern", 70, 600, 122, 20)
    ppt_blue_tag(slide, "feasible subpattern", 250, 600, 128, 20)
    base.add_ppt_line(slide, (430, 612), (490, 612), color=base.BLACK, width=1.7)
    base.add_ppt_text(slide, "vehicle route", 500, 599, 90, 24, size=10.5, align=PP_ALIGN.LEFT)
    base.add_ppt_line(slide, (620, 612), (680, 612), color=base.BLUE, width=1.5, dashed=True)
    base.add_ppt_text(slide, "same-stop UAV sortie", 690, 599, 140, 24, size=10.5, align=PP_ALIGN.LEFT)
    base.add_ppt_green_label(slide, "q95 screen", 830, 600, 78, 20)
    pptx = ppt_save(prs, "Fig2_solution_framework_ppt_icon")
    return png, pdf, pptx


def pattern(canvas, draw, icons, label, color, x, y, n=3, small=False):
    icon_key = {"blue": "tower_blue", "orange": "tower_orange", "green": "tower_green", "purple": "tower_purple", "gold": "tower_gold", "gray": "tower_gray"}[color]
    text_color = {"blue": (0, 102, 204), "orange": base.ORANGE, "green": base.GREEN, "purple": (100, 45, 160), "gold": (213, 150, 0)}[color]
    for i in range(n):
        base.paste_icon(canvas, icons[icon_key], x + i * 14, y, 13, 26)
    base.draw_line_arrow(draw, (x + 4, y + 29), (x + 35, y + 29), color=base.BLUE, width=1.0, dashed=True, arrow=False)
    base.draw_centered(draw, label, (x - 3, y + 30, 46, 15), size=9.5 if small else 10.5, bold=True, fill=text_color)


def ppt_pattern(slide, icons, label, color, x, y, n=3, small=False):
    icon_key = {"blue": "tower_blue", "orange": "tower_orange", "green": "tower_green", "purple": "tower_purple", "gold": "tower_gold", "gray": "tower_gray"}[color]
    text_color = {"blue": (0, 102, 204), "orange": base.ORANGE, "green": base.GREEN, "purple": (100, 45, 160), "gold": (213, 150, 0), "gray": (140, 140, 140)}[color]
    for i in range(n):
        base.add_ppt_icon(slide, icons[icon_key], x + i * 14, y, 13, 26)
    base.add_ppt_line(slide, (x + 4, y + 29), (x + 35, y + 29), color=base.BLUE, width=1.0, dashed=True, arrow=False)
    base.add_ppt_text(slide, label, x - 3, y + 30, 46, 15, size=9.5 if small else 10.5, bold=True, color=text_color)


def draw_schedule_panel(canvas, draw, icons, x, y, title, mode):
    base.draw_centered(draw, title, (x, y, 260, 24), size=14, bold=True)
    nodes = {
        "S1": (x + 45, y + 78),
        "S2": (x + 205, y + 78),
        "S3": (x + 130, y + 145),
        "S4": (x + 45, y + 218),
        "S5": (x + 205, y + 218),
    }
    if mode in {"initial", "risk"}:
        edges = [("S1", "S2"), ("S2", "S3"), ("S1", "S3"), ("S3", "S4"), ("S3", "S5"), ("S4", "S5")]
    elif mode == "removed":
        edges = [("S1", "S2"), ("S1", "S3"), ("S3", "S4"), ("S3", "S5"), ("S4", "S5")]
    else:
        edges = [("S1", "S2"), ("S2", "S3"), ("S3", "S4"), ("S3", "S5"), ("S4", "S5")]
    for a, b in edges:
        base.draw_line_arrow(draw, nodes[a], nodes[b], color=base.BLACK, width=1.5)
    for s, p in nodes.items():
        base.draw_stop(draw, s, *p)
    if mode == "initial":
        pats = [("R1", "blue", "S1"), ("R2", "orange", "S2"), ("R3", "green", "S3"), ("R4", "purple", "S4"), ("R5", "gold", "S5")]
    elif mode == "removed":
        pats = [("R1", "blue", "S1"), ("R3", "green", "S3"), ("R5", "gold", "S5")]
        dashed_box(draw, x + 72, y + 230, 132, 48, color=(160, 160, 160), width=1.0)
        pattern(canvas, draw, icons, "R4", "purple", x + 92, y + 235, small=True)
        pattern(canvas, draw, icons, "R2", "orange", x + 148, y + 235, small=True)
    elif mode == "risk":
        pats = [("R1", "blue", "S1"), ("R2", "orange", "S2"), ("R3", "green", "S3"), ("R4", "purple", "S4"), ("R5", "gold", "S5")]
        green_tag(draw, "q95 repair", x + 150, y + 260, 80, 18)
    else:
        pats = []
        manual_pats = [
            ("R1", "blue", nodes["S1"][0] - 36, nodes["S1"][1] - 58),
            ("R4", "purple", nodes["S2"][0] + 30, nodes["S2"][1] - 58),
            ("R3", "green", nodes["S3"][0] - 25, nodes["S3"][1] - 62),
            ("R5", "gold", nodes["S5"][0] + 30, nodes["S5"][1] - 58),
            ("R2", "orange", nodes["S5"][0] + 28, nodes["S5"][1] - 12),
        ]
        for label, color, px, py in manual_pats:
            pattern(canvas, draw, icons, label, color, px, py)
        return
    offsets = {"S1": (-36, -58), "S2": (30, -58), "S3": (-25, -62), "S4": (-40, -58), "S5": (30, -58)}
    for label, color, s in pats:
        dx, dy = offsets[s]
        pattern(canvas, draw, icons, label, color, nodes[s][0] + dx, nodes[s][1] + dy)


def ppt_schedule_panel(slide, icons, x, y, title, mode):
    base.add_ppt_text(slide, title, x, y, 260, 24, size=14, bold=True)
    nodes = {
        "S1": (x + 45, y + 78),
        "S2": (x + 205, y + 78),
        "S3": (x + 130, y + 145),
        "S4": (x + 45, y + 218),
        "S5": (x + 205, y + 218),
    }
    if mode in {"initial", "risk"}:
        edges = [("S1", "S2"), ("S2", "S3"), ("S1", "S3"), ("S3", "S4"), ("S3", "S5"), ("S4", "S5")]
    elif mode == "removed":
        edges = [("S1", "S2"), ("S1", "S3"), ("S3", "S4"), ("S3", "S5"), ("S4", "S5")]
    else:
        edges = [("S1", "S2"), ("S2", "S3"), ("S3", "S4"), ("S3", "S5"), ("S4", "S5")]
    for a, b in edges:
        base.add_ppt_line(slide, nodes[a], nodes[b], color=base.BLACK, width=1.5)
    for s, p in nodes.items():
        base.add_ppt_stop(slide, s, *p)
    if mode == "initial":
        pats = [("R1", "blue", "S1"), ("R2", "orange", "S2"), ("R3", "green", "S3"), ("R4", "purple", "S4"), ("R5", "gold", "S5")]
    elif mode == "removed":
        pats = [("R1", "blue", "S1"), ("R3", "green", "S3"), ("R5", "gold", "S5")]
        ppt_dashed_box(slide, x + 72, y + 230, 132, 48, color=(160, 160, 160))
        ppt_pattern(slide, icons, "R4", "purple", x + 92, y + 235, small=True)
        ppt_pattern(slide, icons, "R2", "orange", x + 148, y + 235, small=True)
    elif mode == "risk":
        pats = [("R1", "blue", "S1"), ("R2", "orange", "S2"), ("R3", "green", "S3"), ("R4", "purple", "S4"), ("R5", "gold", "S5")]
        base.add_ppt_green_label(slide, "q95 repair", x + 150, y + 260, 80, 18)
    else:
        manual_pats = [
            ("R1", "blue", nodes["S1"][0] - 36, nodes["S1"][1] - 58),
            ("R4", "purple", nodes["S2"][0] + 30, nodes["S2"][1] - 58),
            ("R3", "green", nodes["S3"][0] - 25, nodes["S3"][1] - 62),
            ("R5", "gold", nodes["S5"][0] + 30, nodes["S5"][1] - 58),
            ("R2", "orange", nodes["S5"][0] + 28, nodes["S5"][1] - 12),
        ]
        for label, color, px, py in manual_pats:
            ppt_pattern(slide, icons, label, color, px, py)
        return
    offsets = {"S1": (-36, -58), "S2": (30, -58), "S3": (-25, -62), "S4": (-40, -58), "S5": (30, -58)}
    for label, color, s in pats:
        dx, dy = offsets[s]
        ppt_pattern(slide, icons, label, color, nodes[s][0] + dx, nodes[s][1] + dy)


def draw_fig3(icons):
    h = 740
    canvas, draw = new_canvas(h)
    panels = [
        (30, 28, "Initial schedule", "initial"),
        (350, 28, "After removal", "removed"),
        (670, 28, "After risk-value reinsertion", "risk"),
        (30, 342, "Initial schedule", "initial"),
        (350, 342, "After removal", "removed"),
        (670, 342, "After random reinsertion", "random"),
    ]
    for x, y, title, mode in panels:
        draw_schedule_panel(canvas, draw, icons, x, y, title, mode)
    for y in [156, 470]:
        base.draw_line_arrow(draw, (305, y), (345, y), color=base.BLACK, width=5)
        base.draw_line_arrow(draw, (625, y), (665, y), color=base.BLACK, width=5)
    ly = 675
    draw.rounded_rectangle([base.sp(45), base.sp(ly), base.sp(915), base.sp(ly + 46)], radius=base.sp(6), fill=base.WHITE, outline=(80, 80, 80), width=base.sp(1.1))
    x = 62
    for label, color in [("R1", "blue"), ("R2", "orange"), ("R3", "green"), ("R4", "purple"), ("R5", "gold")]:
        pattern(canvas, draw, icons, label, color, x, ly + 6, small=True)
        x += 78
    dashed_box(draw, 455, ly + 8, 78, 30, color=(160, 160, 160), width=1.0)
    base.draw_centered(draw, "removed pattern", (458, ly + 13, 72, 22), size=8.8)
    green_tag(draw, "q95 feasible repair", 560, ly + 12, 125, 20)
    base.draw_line_arrow(draw, (710, ly + 24), (770, ly + 24), color=base.BLACK, width=1.6)
    base.draw_left(draw, "vehicle route", (780, ly + 12, 92, 24), size=9.5)
    base.draw_line_arrow(draw, (850, ly + 24), (895, ly + 24), color=base.BLUE, width=1.3, dashed=True)
    base.draw_left(draw, "UAV sortie", (902, ly + 12, 52, 24), size=9.0)
    png, pdf = save_canvas(canvas, "Fig8_operator_mechanism_ppt_icon")
    prs, slide = new_ppt(h)
    for x, y, title, mode in panels:
        ppt_schedule_panel(slide, icons, x, y, title, mode)
    for y in [156, 470]:
        base.add_ppt_line(slide, (305, y), (345, y), color=base.BLACK, width=5)
        base.add_ppt_line(slide, (625, y), (665, y), color=base.BLACK, width=5)
    leg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, base.pt(45), base.pt(ly), base.pt(870), base.pt(46))
    leg.fill.solid()
    leg.fill.fore_color.rgb = base.ppt_rgb(base.WHITE)
    leg.line.color.rgb = base.ppt_rgb((80, 80, 80))
    x = 62
    for label, color in [("R1", "blue"), ("R2", "orange"), ("R3", "green"), ("R4", "purple"), ("R5", "gold")]:
        ppt_pattern(slide, icons, label, color, x, ly + 6, small=True)
        x += 78
    ppt_dashed_box(slide, 455, ly + 8, 78, 30, color=(160, 160, 160))
    base.add_ppt_text(slide, "removed pattern", 458, ly + 13, 72, 22, size=8.8)
    base.add_ppt_green_label(slide, "q95 feasible repair", 560, ly + 12, 125, 20)
    base.add_ppt_line(slide, (710, ly + 24), (770, ly + 24), color=base.BLACK, width=1.6)
    base.add_ppt_text(slide, "vehicle route", 780, ly + 12, 92, 24, size=9.5, align=PP_ALIGN.LEFT)
    base.add_ppt_line(slide, (850, ly + 24), (895, ly + 24), color=base.BLUE, width=1.3, dashed=True)
    base.add_ppt_text(slide, "UAV sortie", 902, ly + 12, 52, 24, size=9.0, align=PP_ALIGN.LEFT)
    pptx = ppt_save(prs, "Fig8_operator_mechanism_ppt_icon")
    return png, pdf, pptx


def install(outputs):
    for src_dir in SOURCE_DIRS:
        fig_dir = src_dir / "figures"
        shutil.copy2(outputs["fig1"][0], fig_dir / "Fig0_network_construction_ai.png")
        shutil.copy2(outputs["fig1"][0], fig_dir / "Fig0_network_construction_ppt_icon.png")
        shutil.copy2(outputs["fig1"][1], fig_dir / "Fig0_network_construction_ppt_icon.pdf")
        shutil.copy2(outputs["fig1"][2], fig_dir / "Fig0_network_construction_ppt_icon_source.pptx")
        shutil.copy2(outputs["fig2"][0], fig_dir / "Fig2_solution_framework_ai.png")
        shutil.copy2(outputs["fig2"][0], fig_dir / "Fig2_solution_framework_ppt_icon.png")
        shutil.copy2(outputs["fig2"][1], fig_dir / "Fig2_solution_framework_ppt_icon.pdf")
        shutil.copy2(outputs["fig2"][2], fig_dir / "Fig2_solution_framework_ppt_icon_source.pptx")
        shutil.copy2(outputs["fig3"][0], fig_dir / "Fig8_operator_mechanism_ai.png")
        shutil.copy2(outputs["fig3"][0], fig_dir / "Fig8_operator_mechanism_ppt_icon.png")
        shutil.copy2(outputs["fig3"][1], fig_dir / "Fig8_operator_mechanism_ppt_icon.pdf")
        shutil.copy2(outputs["fig3"][2], fig_dir / "Fig8_operator_mechanism_ppt_icon_source.pptx")


if __name__ == "__main__":
    icons = extend_icons(base.prepare_icons())
    outputs = {
        "fig1": draw_fig1(icons),
        "fig2": draw_fig2(icons),
        "fig3": draw_fig3(icons),
    }
    install(outputs)
    for key, values in outputs.items():
        print(key, "PNG=", values[0])
        print(key, "PDF=", values[1])
        print(key, "PPTX=", values[2])
